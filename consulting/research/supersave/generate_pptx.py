#!/usr/bin/env python3
"""슈퍼세이브 곰보배추시럽 마켓 데스크 리서치 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

def chart_path(name):
    return os.path.join(BASE_DIR, name)

def add_image_safe(slide, path, left, top, width=None, height=None):
    """이미지 파일이 있으면 삽입, 없으면 스킵"""
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top))
        return True
    return False

# ── 색상 정의 ──
ACCENT = RGBColor(0x2d, 0x3a, 0x8c)
ACCENT_LIGHT = RGBColor(0xee, 0xf0, 0xff)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x1a, 0x1a, 0x1a)
MUTED = RGBColor(0x6b, 0x72, 0x80)
GRAY50 = RGBColor(0xf9, 0xfa, 0xfb)
GRAY100 = RGBColor(0xf3, 0xf4, 0xf6)
GRAY200 = RGBColor(0xe5, 0xe7, 0xeb)
GRAY600 = RGBColor(0x4b, 0x55, 0x63)
ORANGE = RGBColor(0xe8, 0x49, 0x1d)
BLUE = RGBColor(0x3b, 0x82, 0xf6)
GREEN = RGBColor(0x10, 0xb9, 0x81)
PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
RED = RGBColor(0xef, 0x44, 0x44)
DARK_BG = RGBColor(0x1e, 0x2a, 0x6e)
GREEN_BG = RGBColor(0x06, 0x4e, 0x3b)

# ── 슬라이드 크기: A4 가로 ──
SLIDE_W = Emu(12192000)  # ~338mm → 13.33in
SLIDE_H = Emu(6858000)   # ~190mm → 7.5in

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── 헬퍼 함수 ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide, color1, color2):
    """Solid fallback (pptx gradient is complex)"""
    add_bg(slide, color1)

def txbox(slide, left, top, width, height, text="", font_size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_paragraph(tf, text="", font_size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p

def add_footer(slide, page_num):
    txbox(slide, 5.5, 7.0, 2.5, 0.3, "BRAND RISE", font_size=7, color=MUTED, align=PP_ALIGN.CENTER)
    txbox(slide, 12.0, 7.0, 1.0, 0.3, str(page_num), font_size=8, color=MUTED, align=PP_ALIGN.RIGHT)

def add_footer_light(slide, page_num):
    txbox(slide, 5.5, 7.0, 2.5, 0.3, "BRAND RISE", font_size=7, color=RGBColor(0xff,0xff,0xff), align=PP_ALIGN.CENTER)
    txbox(slide, 12.0, 7.0, 1.0, 0.3, str(page_num), font_size=8, color=RGBColor(0xff,0xff,0xff), align=PP_ALIGN.RIGHT)

def add_header(slide, section_label, title, page_num):
    txbox(slide, 0.6, 0.3, 4, 0.3, section_label, font_size=8, bold=True, color=ORANGE)
    txbox(slide, 0.6, 0.55, 10, 0.5, title, font_size=20, bold=True, color=BLACK)
    # underline
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.0), Inches(12.0), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    add_footer(slide, page_num)

def add_divider(slide, number, title, desc, page_num, bg_color=DARK_BG):
    add_bg(slide, bg_color)
    txbox(slide, 1.2, 2.0, 3, 1.2, number, font_size=60, bold=True, color=RGBColor(0xff,0xff,0xff))
    # make number semi-transparent via lighter color
    txbox(slide, 1.2, 3.2, 8, 0.8, title, font_size=30, bold=True, color=WHITE)
    txbox(slide, 1.2, 4.2, 8, 0.5, desc, font_size=12, color=RGBColor(0xbb,0xbb,0xbb))
    add_footer_light(slide, page_num)

def add_card_box(slide, left, top, width, height, fill_color=GRAY50):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Pt(3), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ═══════════════════════════════════════════════════
# SLIDE 01: COVER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, DARK_BG, ACCENT)

txbox(slide, 1.2, 1.5, 4, 0.3, "MARKET DESK RESEARCH", font_size=9, bold=True, color=ORANGE)
txbox(slide, 1.2, 2.2, 8, 1.2, "슈퍼세이브\n곰보배추시럽", font_size=36, bold=True, color=WHITE)
tf = txbox(slide, 1.2, 3.8, 8, 0.8, "국내외 시장 환경 분석 및 경쟁 브랜드 벤치마킹", font_size=14, color=RGBColor(0xcc,0xcc,0xcc))
add_paragraph(tf, "건강즙/허벌시럽 카테고리 마켓 데스크 리서치", font_size=14, color=RGBColor(0xcc,0xcc,0xcc))
tf2 = txbox(slide, 1.2, 5.2, 4, 0.6, "Prepared by BRAND RISE", font_size=10, color=RGBColor(0x99,0x99,0x99))
add_paragraph(tf2, "2026.04 | Confidential", font_size=10, color=RGBColor(0x99,0x99,0x99))

# ═══════════════════════════════════════════════════
# SLIDE 02: TOC
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Contents", "목차", "02")

toc_items = [
    ("01", "시장 개요", "시장 정의, 규모, 성장률, 핵심 트렌드", ACCENT),
    ("02", "소비자 분석", "타겟 세그먼트 4개, 구매 동기, 구매 여정", ORANGE),
    ("03", "경쟁사 분석", "국내 9개 브랜드, 포지셔닝 맵, 비교 분석표, 개별 프로필 6건", BLUE),
    ("04", "슈퍼세이브 브랜드 진단", "기업 정보, 제품, 유통, SNS, 브랜딩 현황", GREEN),
    ("05", "유통 채널 분석", "온/오프라인 채널 현황 및 전략 기회", PURPLE),
    ("06", "인사이트 & 기회", "SWOT 분석, 핵심 기회, 전략 방향", RED),
    ("07", "해외 시장 분석", "글로벌 허벌 시장, 벤치마킹 6개 브랜드, 진출 로드맵", RGBColor(0x0e,0xa5,0xe9)),
]

for i, (num, title, desc, color) in enumerate(toc_items):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = 0.8 + col * 6.2
    y = 1.5 + row * 1.5
    add_accent_line(slide, x, y, 0.9, color)
    txbox(slide, x + 0.2, y + 0.05, 5, 0.35, f"{num}. {title}", font_size=14, bold=True, color=BLACK)
    txbox(slide, x + 0.2, y + 0.45, 5, 0.4, desc, font_size=10, color=GRAY600)


# ═══════════════════════════════════════════════════
# SLIDE 03: 01 DIVIDER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_divider(slide, "01", "시장 개요", "Market Overview — 건강즙/곰보배추시럽 시장 정의, 규모, 핵심 트렌드", "03")


# ═══════════════════════════════════════════════════
# SLIDE 04: 시장 정의 & 규모
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "01. 시장 개요", "시장 정의 & 규모", "04")

kpis = [("~6조원", "건강기능식품 전체 (2024)"), ("1.2~1.5조", "건강즙 시장 (추정)"), ("50~150억", "곰보배추시럽 니치 (추정)"), ("4.8%", "건기식 CAGR (5년)")]
for i, (val, label) in enumerate(kpis):
    x = 0.8 + i * 3.0
    add_card_box(slide, x, 1.3, 2.7, 1.0)
    txbox(slide, x + 0.1, 1.35, 2.5, 0.5, val, font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txbox(slide, x + 0.1, 1.85, 2.5, 0.4, label, font_size=8, color=MUTED, align=PP_ALIGN.CENTER)

add_card_box(slide, 0.8, 2.6, 5.8, 2.5)
txbox(slide, 1.0, 2.7, 5.4, 0.3, "시장 정의", font_size=13, bold=True, color=BLACK)
tf = txbox(slide, 1.0, 3.1, 5.4, 1.8, "곰보배추시럽은 단일 카테고리가 아닌 복수 시장의 교차점에 위치합니다.", font_size=10, color=GRAY600)
add_paragraph(tf, "건강즙/건강음료 시장과 전통 민간약초 식품 시장의 교집합이며,", font_size=10, color=GRAY600, space_before=4)
add_paragraph(tf, "식약처 인증 건강기능식품이 아닌 '일반식품'으로 분류되어 기능성 표시가 불가합니다.", font_size=10, color=GRAY600, space_before=4)

add_card_box(slide, 6.9, 2.6, 5.8, 2.5)
txbox(slide, 7.1, 2.7, 5.4, 0.3, "카테고리 분류", font_size=13, bold=True, color=BLACK)
cats = "• 건강즙/건강음료 (핵심)\n• 전통 민간약초 (핵심)\n• 호흡기 건강식품\n• 유아 건강식품 (신흥)"
txbox(slide, 7.1, 3.1, 5.4, 1.5, cats, font_size=10, color=GRAY600)

txbox(slide, 0.8, 5.3, 12, 0.3, "출처: 한국건강기능식품협회, e-나라지표, 헤럴드경제, HIZ 추정", font_size=7, color=MUTED)


# ═══════════════════════════════════════════════════
# SLIDE 05: 시장 규모 추이
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "01. 시장 개요", "건강기능식품 시장 규모 추이", "05")

years = [("2021", "5조 6,902억", "-"), ("2022", "6조 1,498억", "+8.1%"), ("2023", "6조 1,415억", "-0.1%"), ("2024", "5조 9,531억", "-3.1%"), ("2025", "5조 9,626억", "+0.2%")]
for i, (yr, size, growth) in enumerate(years):
    x = 0.8 + i * 2.4
    add_card_box(slide, x, 1.4, 2.2, 1.4)
    txbox(slide, x + 0.1, 1.45, 2.0, 0.3, yr, font_size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txbox(slide, x + 0.1, 1.8, 2.0, 0.4, size, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    g_color = GREEN if growth.startswith("+") else RED if growth.startswith("-") else MUTED
    txbox(slide, x + 0.1, 2.3, 2.0, 0.3, growth, font_size=10, bold=True, color=g_color, align=PP_ALIGN.CENTER)

# 시장 규모 차트 이미지
add_image_safe(slide, chart_path('chart_market_size.png'), 1.0, 3.0, width=11.5)


# ═══════════════════════════════════════════════════
# SLIDE 06: 핵심 트렌드
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "01. 시장 개요", "핵심 트렌드 (2024-2026)", "06")

trends = [
    ("면역력 관심 폭증", "코로나19 이후 호흡기 건강 관심 지속. 면역력 강화 건강식품 수요 상승세", BLUE),
    ("MZ세대 진입", "2030 젊은 세대까지 건강즙 소비층 확대. SNS/숏폼 콘텐츠 중심 구매 여정", PURPLE),
    ("프리미엄화", "저가 경쟁 → 원료 품질·브랜딩으로 경쟁축 이동. 스토리텔링 중요도 상승", GREEN),
    ("유아 건강식품", "유아/어린이 대상 건강즙 시장 급성장. 엄마 소비자 중심 구전 마케팅", ORANGE),
]
for i, (title, desc, color) in enumerate(trends):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.4 + row * 2.5
    add_card_box(slide, x, y, 5.8, 2.0)
    add_accent_line(slide, x, y, 2.0, color)
    txbox(slide, x + 0.3, y + 0.15, 5.2, 0.35, title, font_size=13, bold=True, color=BLACK)
    txbox(slide, x + 0.3, y + 0.6, 5.2, 1.2, desc, font_size=10, color=GRAY600)


# ═══════════════════════════════════════════════════
# SLIDE 07: 성장 동인 & 위험요인
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "01. 시장 개요", "성장 동인 & 위험요인", "07")

# Growth drivers
add_card_box(slide, 0.8, 1.4, 5.8, 4.5)
txbox(slide, 1.0, 1.5, 5.4, 0.35, "성장 동인", font_size=14, bold=True, color=GREEN)
drivers = [
    "• 호흡기 건강 관심 지속 (엔데믹 이후에도 유지)",
    "• 건강즙 소비자 연령층 확대 (중년 → 2030)",
    "• TV홈쇼핑/라이브커머스 채널 성장",
    "• 프리미엄 건강식품 트렌드",
    "• 유아/어린이 건강식품 시장 급성장",
]
tf = txbox(slide, 1.0, 2.0, 5.4, 3.5, drivers[0], font_size=10, color=GRAY600)
for d in drivers[1:]:
    add_paragraph(tf, d, font_size=10, color=GRAY600, space_before=8)

# Risk factors
add_card_box(slide, 6.9, 1.4, 5.8, 4.5)
txbox(slide, 7.1, 1.5, 5.4, 0.35, "위험요인", font_size=14, bold=True, color=RED)
risks = [
    "• 건기식 전체 시장 정체/역성장 (-3.1%, 2024)",
    "• 일반식품 분류 → 기능성 표시 불가",
    "• 영세 업체 난립으로 품질 신뢰 저하",
    "• 원료(곰보배추) 수급 불안정성",
    "• 대형 식품기업 진출 시 경쟁 심화",
]
tf = txbox(slide, 7.1, 2.0, 5.4, 3.5, risks[0], font_size=10, color=GRAY600)
for r in risks[1:]:
    add_paragraph(tf, r, font_size=10, color=GRAY600, space_before=8)


# ═══════════════════════════════════════════════════
# SLIDE 08: 02 DIVIDER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_divider(slide, "02", "소비자 분석", "Consumer Analysis — 타겟 세그먼트, 구매 동기, 구매 여정", "08")


# ═══════════════════════════════════════════════════
# SLIDE 09: 타겟 세그먼트
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "02. 소비자 분석", "타겟 세그먼트 분석", "09")

segments = [
    ("건강관심 중년층", "50~60대", "기존 건강즙 주 소비자\n전통 약재 신뢰도 높음\n오프라인/TV홈쇼핑 선호", ACCENT),
    ("육아맘", "30~40대", "자녀 호흡기 건강 관심\n구전/후기 의존도 높음\n온라인/커뮤니티 활발", ORANGE),
    ("MZ 건강족", "20~30대", "SNS 기반 건강 트렌드\n비주얼/브랜딩 중시\n구독/정기배송 선호", PURPLE),
    ("선물 구매자", "전 연령", "명절/효도 선물 수요\n패키징/세트 구성 중시\n가격보다 신뢰/품질", GREEN),
]
# 소비자 세그먼트 차트
add_image_safe(slide, chart_path('chart_consumer_segments.png'), 0.8, 1.3, width=12.0)


# ═══════════════════════════════════════════════════
# SLIDE 10: 구매 동기 & 여정
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "02. 소비자 분석", "구매 동기 & 구매 여정", "10")

# 구매 여정 차트
add_image_safe(slide, chart_path('chart_purchase_journey.png'), 0.8, 1.3, width=12.0)


# ═══════════════════════════════════════════════════
# SLIDE 11: 03 DIVIDER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_divider(slide, "03", "경쟁사 분석", "Competitor Analysis — Tier별 분류, 비교 분석표, 포지셔닝 맵", "11")


# ═══════════════════════════════════════════════════
# SLIDE 12: 경쟁 구도 개요
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "03. 경쟁사 분석", "경쟁 구도 개요 (Tier별 분류)", "12")

tiers = [
    ("TIER 1 — 브랜딩 선도", ACCENT, [
        ("슈퍼세이브", "(주)승만 · 경기 고양\n리뷰 1,564개 (4.9점)\n52,000원 (선물세트)\nCJ온스타일, 자사몰"),
        ("슬리아 곰백보감", "슬리아 · 3개국 특허\n34,900원 (20포)\n의사/약사 개발 신뢰\n쿠팡, 자사몰"),
    ]),
    ("TIER 2 — 가성비/볼륨", BLUE, [
        ("한국유기농", "(주)한국유기농\n구매자 49만명 · 평점 5.0\n24,900원 (14포)\n쿠마마켓"),
        ("오늘과일", "유아/아이 타겟 특화\n19,980원 (14포)\n국내산 프리미엄\n오아시스, 11번가"),
    ]),
    ("TIER 3 — 농가 직영/로컬", GREEN, [
        ("곰배랑효소랑", "영농조합 · 여주\n26종 라인업\n직접 재배→가공 일관\n자사몰"),
        ("사온데 외", "농가 직영 · 로컬\n원산지 프리미엄\n브랜딩/마케팅 미약\n자사몰, 직거래"),
    ]),
]

for i, (tier_name, color, brands) in enumerate(tiers):
    x = 0.8 + i * 4.2
    txbox(slide, x, 1.4, 3.8, 0.3, tier_name, font_size=9, bold=True, color=color)
    for j, (name, desc) in enumerate(brands):
        y = 1.9 + j * 2.4
        add_card_box(slide, x, y, 3.8, 2.1)
        txbox(slide, x + 0.15, y + 0.1, 3.5, 0.3, name, font_size=12, bold=True, color=BLACK)
        txbox(slide, x + 0.15, y + 0.5, 3.5, 1.5, desc, font_size=9, color=GRAY600)


# ═══════════════════════════════════════════════════
# SLIDE 13: 경쟁사 비교 분석표
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "03. 경쟁사 분석", "경쟁사 비교 분석표", "13")

# Table
rows_data = [
    ["브랜드", "포지셔닝", "스틱 단가", "타겟", "핵심 강점", "핵심 약점", "주요 채널"],
    ["슈퍼세이브", "온 가족 건강방패", "~1,150원", "엄마/가족", "원조, 리뷰4.9", "SNS 취약", "자사몰, CJ"],
    ["슬리아", "의사/약사 개발", "~1,750원", "건강관심층", "3개국 특허", "높은 가격", "쿠팡, 자사몰"],
    ["한국유기농", "공장도 가성비", "~1,780원", "가성비층", "49만 구매자", "브랜딩 약함", "쿠마마켓"],
    ["오늘과일", "유아 프리미엄", "~1,430원", "영유아 엄마", "유아 특화", "인지도 낮음", "오아시스"],
    ["곰배랑효소랑", "농장 직영 신뢰", "~900원", "약초 선호층", "26종, 직접재배", "마케팅 미약", "자사몰"],
    ["새잎", "복합 성분", "미확인", "건강관심층", "삼백초+엘더베리", "인지도 없음", "온라인"],
]

table_shape = slide.shapes.add_table(len(rows_data), 7, Inches(0.6), Inches(1.3), Inches(12.2), Inches(4.5))
table = table_shape.table

for col_idx in range(7):
    cell = table.cell(0, col_idx)
    cell.text = rows_data[0][col_idx]
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(9)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.font.name = "맑은 고딕"
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

for row_idx in range(1, len(rows_data)):
    for col_idx in range(7):
        cell = table.cell(row_idx, col_idx)
        cell.text = rows_data[row_idx][col_idx]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.color.rgb = BLACK
            paragraph.font.name = "맑은 고딕"
        if row_idx == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_LIGHT
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRAY50


# ═══════════════════════════════════════════════════
# SLIDE 14: 포지셔닝 맵
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "03. 경쟁사 분석", "포지셔닝 맵", "14")

# 포지셔닝 맵 차트 이미지
add_image_safe(slide, chart_path('chart_positioning_map.png'), 0.8, 1.2, width=12.0)


# ═══════════════════════════════════════════════════
# SLIDES 15-20: 경쟁사 개별 브랜드 프로필
# ═══════════════════════════════════════════════════
competitors = [
    {
        "name": "슬리아 (곰백보관)",
        "tier": "TIER 1 — 브랜딩 선도", "tier_color": ACCENT,
        "company": "슬리아",
        "positioning": '"30년 경력 의사/약사가 만든 곰보배추시럽"',
        "strength": "3개국 호흡기 특허 성분 '노스릴리브' 보유. 의료진 개발 신뢰도",
        "price": "1박스(20포) 34,900원 / 9박스 216,000원\n스틱 단가 ~1,750원/포 (프리미엄)",
        "channel": "쿠팡, 자사몰, 다나와",
        "feature": "의사·약사 개발이라는 전문가 신뢰 포지셔닝이 명확함. 3개국 특허 성분을 전면에 내세워 경쟁사 대비 과학적 근거가 강함. 프리미엄 가격 정책을 유지하면서도 대량 구매 시 할인으로 재구매 유도.",
        "verdict": "슈퍼세이브 대비 과학적 신뢰도(특허)에서 우위. 단, 가격대가 높아 대중 접근성은 낮음. TV홈쇼핑 미진출.",
        "products": [("곰보배추시럽 15g×20포", "34,900원"), ("9박스 세트", "216,000원")],
        "links": "홈페이지: sliah.com",
        "page": "15",
    },
    {
        "name": "한국유기농",
        "tier": "TIER 2 — 가성비/볼륨", "tier_color": BLUE,
        "company": "(주)한국유기농 · 대표 김성중",
        "positioning": '"저온추출 공법, 국내산 곰보배추+배+생강"',
        "strength": "구매자 49만명, 리뷰 976개, 평점 5.0. 공장도 가성비 소구",
        "price": "1박스(14포) 24,900원 (정가 33,900원, 27% 할인)\n스틱 단가 ~1,780원/포",
        "channel": "쿠마마켓 (주력)",
        "feature": "쿠마마켓 단일 채널에서 49만 구매자를 달성한 볼륨 플레이어. 공장도 가격을 내세우며 가성비 소비자층을 공략. 평점 5.0으로 제품 만족도 최상위.",
        "verdict": "볼륨/리뷰 수에서 슈퍼세이브와 직접 경쟁. 단, 채널 다각화(쿠마마켓 편중)와 브랜딩 투자 부재가 약점.",
        "products": [("곰보배추시럽 1박스(14포)", "24,900원"), ("원료: 곰보배추+배+생강", "국내산")],
        "links": "쿠마마켓: cumamarket.co.kr",
        "page": "16",
    },
    {
        "name": "오늘과일",
        "tier": "TIER 2 — 가성비/볼륨", "tier_color": BLUE,
        "company": "미확인",
        "positioning": '"국내산 프리미엄 유아/아이 건강즙"',
        "strength": "유아/어린이 타겟 특화 포지셔닝. 2만 원 이하 저가 진입으로 접근성 확보",
        "price": "14포 19,980원\n스틱 단가 ~1,430원/포",
        "channel": "오아시스마켓, 11번가",
        "feature": "곰보배추시럽 시장에서 유일하게 유아/아이 전용으로 포지셔닝. 오아시스마켓 등 프리미엄 건강식품 채널에 입점하여 타겟 소비자(건강 관심 높은 엄마)와의 접점 확보.",
        "verdict": "슈퍼세이브가 '온 가족'이라면, 오늘과일은 유아 특화로 세그먼트 차별화. 니치 시장에서 성장 가능성.",
        "products": [("곰보배추시럽 스틱 14포", "19,980원"), ("원료", "국내산 프리미엄")],
        "links": "오아시스마켓 내 검색",
        "page": "17",
    },
    {
        "name": "곰배랑효소랑",
        "tier": "TIER 3 — 농가 직영/로컬", "tier_color": GREEN,
        "company": "곰배랑효소랑 영농조합법인 · 대표 이흥민 · 경기 여주",
        "positioning": '"경기도 여주 곰보배추 전문 농장" — 직접 재배→가공 일관 체계',
        "strength": "곰보배추 전 품목 26종 풀라인업. 직접 재배로 원재료 신뢰도 최상",
        "price": "달임액 30포 27,000원 / 차 70포 53,000원 / 환 200g 25,000원\n스틱 단가 ~900원/포 (최저가)",
        "channel": "자사몰 (gombobaechu.com)",
        "feature": "곰보배추 전문 농장 직영으로 원재료부터 최종 제품까지 일관된 품질 관리. 26종이라는 압도적 라인업이 강점이나, 온라인 마케팅과 브랜딩 투자가 매우 미약.",
        "verdict": "원재료 신뢰도에서 최상위이나 브랜딩/마케팅 전무. 원료 파트너로서의 가치 존재.",
        "products": [("달임액 30포", "27,000원"), ("곰보배추차 70포", "53,000원"), ("곰보배추환 200g", "25,000원"), ("외 23종", "다양")],
        "links": "홈페이지: gombobaechu.com",
        "page": "18",
    },
    {
        "name": "사온데",
        "tier": "TIER 3 — 농가 직영/로컬", "tier_color": GREEN,
        "company": "사온데 · 낙동강 상수원보호구역",
        "positioning": '"낙동강 상수원보호구역 재배" — 원산지 프리미엄',
        "strength": "상수원보호구역이라는 청정 원산지 스토리. 효소/발효액/시럽/건초/즙 등 다양한 형태",
        "price": "다양 (제품 형태별 상이)",
        "channel": "자사몰 (saonde.com)",
        "feature": "낙동강 상수원보호구역이라는 독보적인 원산지 스토리가 핵심 자산. 다양한 형태로 제품화. 농가 직영 특유의 진정성과 원료 신뢰도가 높으나, 마케팅/브랜딩 활동은 미약.",
        "verdict": "원산지 스토리텔링 잠재력이 높으나 현재 활용 못하는 상태. 브랜딩 지원 시 급성장 가능성.",
        "products": [("곰보배추 효소", "별도 확인"), ("곰보배추 시럽", "별도 확인")],
        "links": "홈페이지: saonde.com",
        "page": "19",
    },
    {
        "name": "새잎",
        "tier": "TIER 2 — 가성비/볼륨", "tier_color": BLUE,
        "company": "미확인",
        "positioning": '"곰보배추 + 삼백초 + 엘더베리 복합"',
        "strength": "복합 성분 차별화 — 곰보배추 단독이 아닌 삼백초·엘더베리 블렌딩",
        "price": "45,100원 (스틱)",
        "channel": "온라인",
        "feature": "곰보배추 시럽 시장에서 복합 성분 전략을 취하는 유일한 브랜드. 가격대가 높아 프리미엄을 지향하나, 브랜드 인지도가 거의 없어 소비자 접근에 한계.",
        "verdict": "복합 성분 차별화 전략은 슈퍼세이브 라인업 확장 시 참고 가능. 현재 직접 위협은 낮음.",
        "products": [("곰보배추시럽 삼백초 엘더베리진액", "45,100원")],
        "links": "온라인 쇼핑몰 (상세 미확인)",
        "page": "20",
    },
    # ── 국내 추가 10개 ──
    {
        "name": "농부건강농원", "tier": "TIER 3 — 농가 직영/로컬", "tier_color": GREEN,
        "company": "농업회사법인(주) 농부건강농원 · 대표 고재은 · 경북 예천",
        "positioning": '"정직한 곰보배추즙" — 예천 농부 부부가 직접 재배·가공',
        "strength": "경북 예천 자체 농장, 재배→가공 일관 체계, 대추 첨가로 맛 보완",
        "price": "곰보배추즙 100ml×50포 (대추첨가)\n추정 3~5만원대",
        "channel": "사이소 마켓(예천 농특산물), 자사몰(farmerfarm.kr)",
        "feature": "소규모 농가형 건강원. 산지직송 모델. 대추를 첨가하여 곰보배추 쓴맛을 보완. 브랜딩/SNS 마케팅 거의 없는 전형적 농가 직영 모델.",
        "verdict": "슈퍼세이브의 TV홈쇼핑 유통 대비 로컬 직거래에 머물러 직접 위협 낮음. '자체 재배' 신뢰도는 강점.",
        "products": [("정직한 곰보배추즙 100ml×50포", "추정 3~5만원")],
        "links": "홈페이지: farmerfarm.kr", "page": "21",
    },
    {
        "name": "동진농장 (부안쇼핑)", "tier": "TIER 3 — 농가 직영/로컬", "tier_color": GREEN,
        "company": "동진농장영농조합법인 · 대표 양천 · 전북 부안",
        "positioning": '"전북 부안 곰보배추 전문 농가" — 부안 산지 프리미엄',
        "strength": "부안 곰보배추 산지, 다린즙/발효원액/환/분말 다양한 제형, 낮은 가격",
        "price": "다린즙 30개(100ml) 22,000원\n발효원액 1.5L 35,000원",
        "channel": "부안쇼핑(buanshop.kr), 옥션",
        "feature": "제형이 다양하고(다린즙, 발효원액, 환, 분말) 가격이 낮아 가성비 포지셔닝 명확. 브랜딩 투자 최소. 전통 즙/환 중심.",
        "verdict": "가격 경쟁력 우위이나 시럽이 아닌 전통 즙/환 중심. 슈퍼세이브 타겟(유아/가족)과 다른 중장년층 집중.",
        "products": [("곰보배추 다린즙 30개", "22,000원"), ("발효원액 1.5L", "35,000원")],
        "links": "홈페이지: buanshop.kr", "page": "22",
    },
    {
        "name": "곰베진 (아름채움)", "tier": "TIER 2 — 가성비/볼륨", "tier_color": BLUE,
        "company": "아름채움",
        "positioning": '"곰보배추 + 삼백초 + 엘더베리 복합 추출물 스틱"',
        "strength": "복수 원료 배합 복합 포뮬러, 스틱 형태 휴대성",
        "price": "곰베진 14포 — 가격 미확인\n(카카오톡 스토어 톡딜)",
        "channel": "카카오톡 스토어 (톡딜)",
        "feature": "곰보배추 단독이 아닌 삼백초·엘더베리 복합 건강 스틱. 카카오톡 소셜커머스 집중. 인지도/SNS 존재감 낮음.",
        "verdict": "복합 원료는 차별점이나 유통(카카오톡 톡딜)이 매우 제한적. 슈퍼세이브 대비 규모/브랜딩 열세.",
        "products": [("곰보배추 삼백초 추출물 스틱 14포", "미확인")],
        "links": "카카오톡 스토어", "page": "23",
    },
    {
        "name": "한우리약초 (동의식품)", "tier": "TIER 3 — 농가 직영/로컬", "tier_color": GREEN,
        "company": "동의식품 · 경기 고양",
        "positioning": '"국내산 건조 약초 전문" — 200종+ 약초/건초 유통',
        "strength": "200종 이상 약초 카탈로그, 건조 원물 중심, DIY 건강식품 소비자 타겟",
        "price": "곰보배추 건조 300g — 미확인",
        "channel": "자사몰 (hanwoore.com)",
        "feature": "원물 판매 모델. 소비자가 직접 달여 먹는 전통적 소비 패턴 대응. 가공식품이 아닌 원재료 유통.",
        "verdict": "직접 경쟁자라기보다 시장 인접 플레이어. 원물 구매 소비자는 슈퍼세이브의 편의형 시럽 타겟과 다름.",
        "products": [("곰보배추 건조 300g", "미확인")],
        "links": "홈페이지: hanwoore.com", "page": "24",
    },
    {
        "name": "GNM자연의품격", "tier": "TIER 1 — 브랜딩 선도", "tier_color": ACCENT,
        "company": "(주)지엔엠라이프 · 공동대표 이건수, 박유영, 유재국",
        "positioning": '"누구나 건강할 자격" — 건강즙 시장 국내 1위급, 9년 연속 소비자만족지수 1위',
        "strength": "2023년 매출 약 970억원, 유기농 양배추즙 누적 1.5억포, 올리브영/SSG/홈쇼핑 전 채널 입점",
        "price": "유기농 양배추즙 30포 15,000~20,000원",
        "channel": "자사몰, 올리브영, SSG, 롯데홈쇼핑, 쿠팡, 코스트코",
        "feature": "건강즙 전 카테고리 커버. 올리브영 입점으로 2030세대 포괄. 건기식까지 사업 확장. 곰보배추 카테고리 미진출이나 진입 시 최대 위협.",
        "verdict": "매출 규모 압도적(970억). 곰보배추 미진출이 슈퍼세이브의 기회 — 카테고리 선점이 핵심.",
        "products": [("유기농 양배추즙 30포", "~20,000원"), ("흑마늘진액 30포", "~18,000원")],
        "links": "홈페이지: gnmart.co.kr", "page": "25",
    },
    {
        "name": "참앤들황토농원", "tier": "TIER 1 — 브랜딩 선도", "tier_color": ACCENT,
        "company": "(주)참앤들황토농원 · 대표 권종성 · 부산 강서구",
        "positioning": '"NFC 착즙 건강즙 전문" — 직접 생산, 중간 유통 제거',
        "strength": "2024년 매출 약 83억원, NFC 착즙 공법, GS홈쇼핑 입점",
        "price": "유기농 양배추즙 30포\n15,000~30,000원대",
        "channel": "자사몰(htfarm.co.kr), GS홈쇼핑, SSG, 네이버",
        "feature": "NFC 100% 착즙 강조. GS홈쇼핑 브랜드관 운영. 중간 유통 제거로 합리적 가격. 곰보배추 미진출이나 다품종 전략.",
        "verdict": "GS홈쇼핑 입점으로 슈퍼세이브(CJ온스타일)와 유사 유통 전략. 곰보배추 진입 가능성 존재.",
        "products": [("양배추즙 30포", "~20,000원"), ("칡즙/흑염소 30포", "~30,000원")],
        "links": "홈페이지: htfarm.co.kr", "page": "26",
    },
    {
        "name": "김재식헬스푸드", "tier": "TIER 2 — 가성비/볼륨", "tier_color": BLUE,
        "company": "농업회사법인 김재식헬스푸드(주) · 대표 김재식 · 경북 영천",
        "positioning": '"진짜 양배추즙" — 저온착즙 100%, 22년 업력 농가 기반',
        "strength": "2003년 설립(22년), 매출 약 36억원(2023), GS홈쇼핑/CJ온스타일 입점",
        "price": "양배추즙 100ml×30포 21,900원\n(정가 40,000원, ~45% 할인)",
        "channel": "자사몰(drkims.co.kr), GS홈쇼핑, G마켓",
        "feature": "정가 대비 40~50% 할인 가격 전략. 양배추즙 중심. 종합 농산 가공 기업(즙/잼/면류/포도주). 곰보배추 미진출.",
        "verdict": "농가+TV홈쇼핑 모델이 슈퍼세이브와 유사. 브랜딩 수준은 슈퍼세이브보다 낮음.",
        "products": [("진짜 양배추즙 30포", "21,900원"), ("ABC주스 30포", "22,900원")],
        "links": "홈페이지: drkims.co.kr", "page": "27",
    },
    {
        "name": "유기농마루", "tier": "TIER 1 — 브랜딩 선도", "tier_color": ACCENT,
        "company": "농업회사법인 (주)유기농마루 · 대표 김상범 · 광주 광산구",
        "positioning": '"원료까지 생각하는 친환경 브랜드" — 유기농 인증 건강즙 전문',
        "strength": "2022년 매출 약 132억원, 유기농 인증 원료 특화, SSG 브랜드관 운영",
        "price": "유기농 양배추즙 30포\n15,000~20,000원대",
        "channel": "자사몰(organicmaru.co.kr), SSG, G마켓, 네이버",
        "feature": "'유기농'이 핵심 정체성. 양배추즙 중심 + 석류즙/매실청 확장. 곰보배추 미진출이나 시장 성장 시 진입 가능.",
        "verdict": "유기농 인증은 슈퍼세이브 미보유 강점. 곰보배추 미진출이나 유기농 곰보배추즙 진입 가능성.",
        "products": [("유기농 양배추즙 30포", "~20,000원"), ("석류즙/매실청", "다양")],
        "links": "홈페이지: organicmaru.co.kr", "page": "28",
    },
    {
        "name": "천호엔케어", "tier": "TIER 1 — 브랜딩 선도", "tier_color": ACCENT,
        "company": "(주)천호엔케어 · 1984년 창업 · 40년 업력",
        "positioning": '"한방건강식품 선도 기업" — 달팽이엑기스로 시장 개척',
        "strength": "2023년 매출 414억원, 74% 반복구매율, TV홈쇼핑 전 채널, 일본/미국 수출",
        "price": "프리미엄 (제품별 상이)",
        "channel": "자사몰, CJ/GS/롯데/현대 전 홈쇼핑, 온라인 전 채널",
        "feature": "건강즙 시장 원조. 달팽이엑기스/강화사자발쑥 히트상품. 한방약초 기반이라 곰보배추 진입 가능성 높음. 74% 반복구매율.",
        "verdict": "규모(414억)와 채널 파워 압도적. 곰보배추 미출시이나 잠재적 고위협. 카테고리 선점이 관건.",
        "products": [("달팽이엑기스", "프리미엄"), ("강화사자발쑥", "프리미엄")],
        "links": "홈페이지: chunhomall.com", "page": "29",
    },
    {
        "name": "순수식품", "tier": "TIER 2 — 가성비/볼륨", "tier_color": BLUE,
        "company": "농업회사법인 (주)순수식품 · 대표 김인섭 · 진주/서울",
        "positioning": '"싸다좋다 순수식품" — 가격 경쟁력 + 셀프 메디케이션',
        "strength": "착즙/저온추출/동결건조 다양한 기술, 엘더베리시럽 보유, 매출 약 22억(2023)",
        "price": "도라지배즙 80ml×100포\n~25,000~35,000원",
        "channel": "자사몰(soonsoofood.co.kr), SSG, 네이버",
        "feature": "'싸다좋다' 슬로건. 가격 경쟁력 핵심. 도라지배즙~콜라겐젤리까지 넓은 라인업. 엘더베리시럽 보유.",
        "verdict": "매출 규모 유사(22억). 가격경쟁+다품종 vs 슈퍼세이브 곰보배추 특화+프리미엄. 곰보배추 진입 시 가격 경쟁 발생 가능.",
        "products": [("도라지배즙 100포", "~30,000원"), ("엘더베리시럽", "미확인")],
        "links": "홈페이지: soonsoofood.co.kr", "page": "30",
    },
]

# ── 해외 벤치마킹 10개 ──
overseas_brands = [
    {
        "name": "Ricola", "country": "스위스", "established": "1930년",
        "revenue": "CHF 3.39억+ (2018), 2023 역대 최고 매출",
        "products": [("13가지 스위스 알프스 허브 블렌드 캔디/시럽", "60종+")],
        "positioning": "스위스 알프스 허브의 힘으로 목과 호흡기를 케어하는 프리미엄 허벌 캔디",
        "feature": "변하지 않은 13허브 레시피로 원조 포지셔닝. Swiss Made 원산지 프리미엄. 2023 B Corp 인증. 50개국+ 수출, 연간 70억개 생산.",
        "benchmark": "곰보배추 단일 원료를 리콜라의 시그니처 레시피처럼 원조 포지셔닝으로 승화. 원산지 스토리 활용.",
        "links": "ricola.com", "page": "31",
    },
    {
        "name": "Sambucol", "country": "이스라엘", "established": "1991년",
        "revenue": "비공개 (60개국 유통, 미국 엘더베리 시럽 1위)",
        "products": [("블랙 엘더베리 시럽", "오리지널/키즈/이뮨 포르테")],
        "positioning": "바이러스학자가 개발한 과학 기반 엘더베리 면역 시럽",
        "feature": "Dr. Mumcuoglu가 항바이러스 성분 학술 규명 후 상품화. 신종플루/독감 시즌마다 매출 급증. 면역=엘더베리=Sambucol 공식.",
        "benchmark": "곰보배추 호흡기 효능의 임상/학술 데이터 확보 → '연구 기반 허벌 시럽' 포지셔닝.",
        "links": "sambucolusa.com", "page": "32",
    },
    {
        "name": "Nin Jiom Pei Pa Koa (배파고)", "country": "홍콩", "established": "1946년 (레시피 청나라)",
        "revenue": "HKD 3.5억 (약 600억원, 2014)",
        "products": [("배파고 — 15가지 한방 허브+꿀 목 시럽", "20개국+")],
        "positioning": "효심에서 탄생한 150년 전통 한방 목 시럽",
        "feature": "'어머니를 그리며'라는 효도 스토리 내러티브. 브로드웨이/셀럽 자발적 바이럴. 캔디/차 등 포맷 확장으로 젊은층 공략.",
        "benchmark": "곰보배추의 '할머니 손 레시피' 전통 스토리텔링 + 캔디/차 서브 포맷 확장.",
        "links": "ninjiom.com", "page": "33",
    },
    {
        "name": "Gaia Herbs", "country": "미국", "established": "1987년",
        "revenue": "약 $82M (2026 추정), 리브랜딩 후 4년간 3배 성장",
        "products": [("액상 허브 추출물 60종+", "Bronchial Wellness 등")],
        "positioning": "농장에서 소비자까지 — 세계 최초 허브 추적 가능 프리미엄 허벌 보충제",
        "feature": "Meet Your Herbs 원산지 추적 플랫폼. DNA 검증 라벨. 자체 유기농 농장. 투명성만으로 4년 3배 매출 성장.",
        "benchmark": "곰보배추 원산지 재배→수확→가공 추적 시스템 구축 → D2C 자사몰 전략의 교과서.",
        "links": "gaiaherbs.com", "page": "34",
    },
    {
        "name": "Comvita", "country": "뉴질랜드", "established": "1974년",
        "revenue": "NZD 2.04억 (FY2024), NZX 상장사",
        "products": [("UMF 마누카 꿀 / 프로폴리스 시럽", "면역 로젠지")],
        "positioning": "뉴질랜드 원시 자연의 과학적 벌꿀 기반 면역 솔루션",
        "feature": "UMF 등급 시스템 공동 개발로 산업 표준 정립. 블록체인 원산지 추적. 로젠지/스낵 확장 중.",
        "benchmark": "곰보배추 자체 품질 등급 체계를 만들면 카테고리 리더 가능.",
        "links": "comvita.com", "page": "35",
    },
    {
        "name": "AG1", "country": "미국", "established": "2009년",
        "revenue": "약 $600M (2024 추정), 기업가치 $1.2B",
        "products": [("AG1 그린 파우더 — 단일 SKU", "75종 영양소 올인원")],
        "positioning": "하루 한 스쿱으로 끝내는 올인원 영양 솔루션",
        "feature": "단일 제품 $600M 매출. 구독 모델 핵심. 팟캐스트 광고주 3위. 공항/스타벅스 오프라인 확장 중.",
        "benchmark": "곰보배추시럽을 단일 히어로 SKU+정기구독 모델로 설계, 인플루언서 파트너십 집중.",
        "links": "drinkag1.com", "page": "36",
    },
    {
        "name": "Traditional Medicinals", "country": "미국", "established": "1974년",
        "revenue": "$170M (2024), 전년비 15% 성장",
        "products": [("약용 허벌 티 60종+", "소화/면역/수면")],
        "positioning": "약초학자가 만든 미국 No.1 메디컬 그레이드 허벌 티",
        "feature": "약초학자 Rosemary Gladstar 창업. 연간 7.82억 티백. 유기농+공정무역 이중 가치.",
        "benchmark": "한방/약초 전문가 권위 콘텐츠 + 시럽 외 티백/분말 포맷 확장.",
        "links": "traditionalmedicinals.com", "page": "37",
    },
    {
        "name": "Pukka Herbs", "country": "영국", "established": "2001년",
        "revenue": "GBP 35.9M (2017 Unilever 인수), 연 25~28% 성장",
        "products": [("프리미엄 유기농 허벌 티", "Elderberry & Echinacea 등")],
        "positioning": "사람, 식물, 지구의 건강을 연결하는 유기농 티 브랜드",
        "feature": "B Corp + 1% for the Planet. 미션 드리븐 브랜딩. 45개국 수출. 2017년 Unilever 인수.",
        "benchmark": "브랜드 미션/가치 체계를 먼저 확립, 모든 터치포인트에 일관 반영하는 미션 드리븐 브랜딩.",
        "links": "pukkaherbs.com", "page": "38",
    },
    {
        "name": "Herbion Naturals", "country": "파키스탄", "established": "1983년",
        "revenue": "비공개 (40개국+, 250종+ 제품)",
        "products": [("Linkus 기침 시럽", "Ivy Leaf Cough Drops")],
        "positioning": "40년 전통 허벌 기침 시럽 글로벌 리더",
        "feature": "중앙아시아 거점 → 40개국 확장. GMP 인증 4개 공장. 시럽→스킨케어까지 250종. Target/Walmart/CVS 입점.",
        "benchmark": "동남아/중화권 등 한방 친화적 거점 시장 먼저 공략 → 글로벌 확장 단계적 전략.",
        "links": "pk.herbion.com", "page": "39",
    },
    {
        "name": "Umeken", "country": "일본", "established": "1978년 (뿌리 1958년)",
        "revenue": "비공개 (연간 200톤, 10억알+ 생산)",
        "products": [("우메보시 볼 / 스피루리나 볼", "45종 한방 허브 압축볼")],
        "positioning": "에도시대 약환 전통을 현대 과학으로 재해석한 일본 No.1 허브 보충제",
        "feature": "에도시대 도야마현 약환 전통 헤리티지. 45가지 허브 한알 압축 기술. 자체 스피루리나 농장. 30년+ 해외 운영.",
        "benchmark": "곰보배추 한국 전통 민간요법 역사를 헤리티지 스토리로 구조화 + 환/정제 간편 제형 확장.",
        "links": "umeken.com", "page": "40",
    },
]

for comp in competitors:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.fill.background()
    header_label = f"국내  |  {comp['name']}"
    txbox(slide, 0.6, 0.15, 8, 0.35, header_label, font_size=16, bold=True, color=WHITE)

    # Tier badge
    add_card_box(slide, 10.5, 0.15, 2.5, 0.35, comp["tier_color"])
    txbox(slide, 10.5, 0.17, 2.5, 0.3, comp["tier"], font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left side - info
    info_items = [
        ("운영사", comp["company"]),
        ("포지셔닝", comp["positioning"]),
        ("핵심 강점", comp["strength"]),
        ("가격대", comp["price"]),
        ("유통 채널", comp["channel"]),
    ]
    y_pos = 1.0
    for label, value in info_items:
        txbox(slide, 0.8, y_pos, 1.0, 0.25, f"●  {label}", font_size=9, bold=True, color=BLACK)
        txbox(slide, 2.0, y_pos, 4.8, 0.5, value, font_size=9, color=GRAY600)
        y_pos += 0.55

    # Feature box
    add_card_box(slide, 0.8, y_pos + 0.1, 5.8, 1.8, GRAY50)
    txbox(slide, 1.0, y_pos + 0.2, 5.4, 0.25, "브랜드 특징", font_size=9, bold=True, color=ACCENT)
    txbox(slide, 1.0, y_pos + 0.5, 5.4, 1.3, comp["feature"], font_size=9, color=GRAY600)

    # Right side
    txbox(slide, 7.0, 1.0, 3, 0.25, "온라인 채널", font_size=9, bold=True, color=MUTED)
    txbox(slide, 7.0, 1.3, 5.5, 0.3, comp["links"], font_size=9, color=BLUE)

    txbox(slide, 7.0, 1.8, 3, 0.25, "주요 제품", font_size=9, bold=True, color=MUTED)
    for pi, (pname, pprice) in enumerate(comp["products"][:4]):
        px = 7.0 + (pi % 2) * 2.9
        py = 2.15 + (pi // 2) * 1.0
        add_card_box(slide, px, py, 2.7, 0.85, GRAY50)
        txbox(slide, px + 0.1, py + 0.05, 2.5, 0.3, pname, font_size=8, bold=True, color=BLACK)
        txbox(slide, px + 0.1, py + 0.4, 2.5, 0.25, pprice, font_size=9, bold=True, color=ACCENT)

    # Verdict
    verdict_y = 4.5
    add_accent_line(slide, 7.0, verdict_y, 1.2, ORANGE)
    add_card_box(slide, 7.1, verdict_y, 5.5, 1.2, RGBColor(0xff, 0xf7, 0xed))
    txbox(slide, 7.3, verdict_y + 0.05, 5.1, 0.25, "BRAND RISE 분석", font_size=8, bold=True, color=ORANGE)
    txbox(slide, 7.3, verdict_y + 0.35, 5.1, 0.8, comp["verdict"], font_size=9, color=BLACK)

    # 제품 이미지 영역 (PPT에서 이미지 교체 가능)
    add_card_box(slide, 7.0, 5.9, 5.6, 1.0, GRAY100)

    add_footer(slide, comp["page"])


# ═══════════════════════════════════════════════════
# SLIDES 31-40: 해외 벤치마킹 브랜드 프로필
# ═══════════════════════════════════════════════════
for ob in overseas_brands:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    txbox(slide, 0.6, 0.15, 8, 0.35, f"해외  |  {ob['name']}", font_size=16, bold=True, color=WHITE)

    # Country badge
    add_card_box(slide, 10.5, 0.15, 2.5, 0.35, RGBColor(0xee, 0xf0, 0xff))
    txbox(slide, 10.5, 0.17, 2.5, 0.3, ob["country"], font_size=9, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Left side - info
    info_items = [
        ("설립", ob["established"]),
        ("매출/규모", ob["revenue"]),
        ("포지셔닝", ob["positioning"]),
    ]
    y_pos = 1.0
    for label, value in info_items:
        txbox(slide, 0.8, y_pos, 1.0, 0.25, f"●  {label}", font_size=9, bold=True, color=BLACK)
        txbox(slide, 2.0, y_pos, 4.8, 0.5, value, font_size=9, color=GRAY600)
        y_pos += 0.55

    # Feature box
    add_card_box(slide, 0.8, y_pos + 0.1, 5.8, 1.8, GRAY50)
    txbox(slide, 1.0, y_pos + 0.2, 5.4, 0.25, "핵심 성공 요인", font_size=9, bold=True, color=ACCENT)
    txbox(slide, 1.0, y_pos + 0.5, 5.4, 1.3, ob["feature"], font_size=9, color=GRAY600)

    # Right side
    txbox(slide, 7.0, 1.0, 3, 0.25, "홈페이지", font_size=9, bold=True, color=MUTED)
    txbox(slide, 7.0, 1.3, 5.5, 0.3, ob["links"], font_size=9, color=BLUE)

    txbox(slide, 7.0, 1.8, 3, 0.25, "주요 제품", font_size=9, bold=True, color=MUTED)
    for pi, (pname, pprice) in enumerate(ob["products"][:4]):
        px = 7.0 + (pi % 2) * 2.9
        py = 2.15 + (pi // 2) * 1.0
        add_card_box(slide, px, py, 2.7, 0.85, GRAY50)
        txbox(slide, px + 0.1, py + 0.05, 2.5, 0.3, pname, font_size=8, bold=True, color=BLACK)
        txbox(slide, px + 0.1, py + 0.4, 2.5, 0.25, pprice, font_size=9, bold=True, color=ACCENT)

    # Benchmark point
    verdict_y = 4.5
    add_accent_line(slide, 7.0, verdict_y, 1.2, ACCENT)
    add_card_box(slide, 7.1, verdict_y, 5.5, 1.2, RGBColor(0xee, 0xf0, 0xff))
    txbox(slide, 7.3, verdict_y + 0.05, 5.1, 0.25, "벤치마킹 포인트", font_size=8, bold=True, color=ACCENT)
    txbox(slide, 7.3, verdict_y + 0.35, 5.1, 0.8, ob["benchmark"], font_size=9, color=BLACK)

    # Image placeholder
    add_card_box(slide, 7.0, 5.9, 5.6, 1.0, GRAY100)

    add_footer(slide, ob["page"])


# ═══════════════════════════════════════════════════
# SLIDE 41: 04 DIVIDER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_divider(slide, "04", "슈퍼세이브 브랜드 진단", "Brand Diagnosis — 기업 정보, 제품, 유통, SNS, 브랜딩 현황 점검", "41", GREEN_BG)


# ═══════════════════════════════════════════════════
# SLIDE 22: 기업 정보 & 제품 라인업
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "04. 슈퍼세이브 브랜드 진단", "기업 정보 & 제품 라인업", "42")

add_card_box(slide, 0.8, 1.4, 5.8, 4.8)
add_accent_line(slide, 0.8, 1.4, 4.8, ACCENT)
txbox(slide, 1.1, 1.5, 5.2, 0.3, "COMPANY", font_size=9, bold=True, color=MUTED)
company_info = (
    "법인명: 주식회사 승만\n"
    "대표: 신주리\n"
    "설립: 2022년 (추정)\n"
    "소재: 경기 고양시 일산동구\n"
    "슬로건: \"온 가족 건강방패 슈퍼세이브\"\n"
    "스토리: \"두 아들의 엄마가 만든 우리집 상비식품\"\n"
    "공식몰: supersave.kr (카페24)\n"
    "추정 매출: 연 10~20억원대"
)
txbox(slide, 1.1, 1.9, 5.2, 3.8, company_info, font_size=10, color=GRAY600)

add_card_box(slide, 6.9, 1.4, 5.8, 4.8)
add_accent_line(slide, 6.9, 1.4, 4.8, ORANGE)
txbox(slide, 7.2, 1.5, 5.2, 0.3, "PRODUCT", font_size=9, bold=True, color=MUTED)
product_info = (
    "① 곰보배추시럽 선물세트 (10g×15포×3박스)\n"
    "   정가 86,850원 → 할인가 52,000원 (~40%↓)\n\n"
    "② 곰보배추시럽 300ml 병\n"
    "   오아시스마켓 등에서 판매\n\n"
    "③ 기타 건강식품 라인업 확장 중\n\n"
    "※ 현재 선물세트 SOLD OUT 상태 (2026.04)"
)
txbox(slide, 7.2, 1.9, 5.2, 3.8, product_info, font_size=10, color=GRAY600)

# 가격 비교 차트 (제품 이미지 대신)
add_image_safe(slide, chart_path('chart_price_comparison.png'), 0.8, 6.0, width=5.8, height=1.2)


# ═══════════════════════════════════════════════════
# SLIDE 23: 유통 & SNS 현황
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "04. 슈퍼세이브 브랜드 진단", "유통 & SNS 현황", "43")

add_card_box(slide, 0.8, 1.4, 5.8, 4.8)
txbox(slide, 1.0, 1.5, 5.4, 0.3, "유통 채널", font_size=13, bold=True, color=ACCENT)
channels = (
    "• CJ온스타일 — 단독 런칭, TV홈쇼핑 핵심 채널\n"
    "• 자사몰 (supersave.kr) — 카페24 기반\n"
    "• 오아시스마켓 — 프리미엄 건강식품 채널\n"
    "• G마켓 — 뉴온아이앤씨 명의 판매\n"
    "• 네이버 스마트스토어 — 검색 유입"
)
txbox(slide, 1.0, 2.0, 5.4, 3.5, channels, font_size=10, color=GRAY600)

add_card_box(slide, 6.9, 1.4, 5.8, 4.8)
txbox(slide, 7.1, 1.5, 5.4, 0.3, "SNS 현황", font_size=13, bold=True, color=ACCENT)
sns = (
    "• 인스타그램: @supergombo\n"
    "  → 팔로워/콘텐츠 규모 확인 필요\n"
    "  → 바이오: \"곰보배추시럽 원조, 네이버 1위\"\n\n"
    "• 블로그/카페: 네이버 블로그 후기 다수\n"
    "  → 체험단/협찬 후기 중심\n\n"
    "• 유튜브/틱톡: 미운영 (기회 영역)"
)
txbox(slide, 7.1, 2.0, 5.4, 3.5, sns, font_size=10, color=GRAY600)


# ═══════════════════════════════════════════════════
# SLIDE 24: SWOT (04 브랜드 진단 내)
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "04. 슈퍼세이브 브랜드 진단", "SWOT 분석", "44")

swot = [
    ("S — 강점", BLUE, "• CJ온스타일 입점 (TV홈쇼핑 유일)\n• 리뷰 1,564개 / 평점 4.9\n• \"원조\" 포지셔닝\n• 엄마 소비자 신뢰 스토리"),
    ("W — 약점", ORANGE, "• SNS/디지털 마케팅 취약\n• 네이밍 혼선 (슈퍼세이브 = ?)\n• 제품 라인업 단조\n• 브랜드 아이덴티티 미약"),
    ("O — 기회", GREEN, "• 호흡기 건강 관심 지속\n• MZ세대 건강즙 시장 진입\n• 유아 건강식품 급성장\n• 라이브커머스/숏폼 채널"),
    ("T — 위협", RED, "• 영세 업체 난립 / 품질 신뢰↓\n• 대형 식품기업 진출 가능성\n• 일반식품 분류 (기능성 표시 불가)\n• 원료 수급 불안정"),
]
# SWOT 차트 이미지
add_image_safe(slide, chart_path('chart_swot.png'), 0.8, 1.2, width=12.0)


# ═══════════════════════════════════════════════════
# SLIDE 25: 05 DIVIDER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_divider(slide, "05", "유통 채널 분석", "Distribution Channel — 온/오프라인 채널 현황 및 전략 기회", "45")


# ═══════════════════════════════════════════════════
# SLIDE 26: 유통 채널 비중
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "05. 유통 채널 분석", "채널별 비중 & 특성", "46")

# 채널 비중 파이차트
add_image_safe(slide, chart_path('chart_channel_pie.png'), 0.8, 1.2, width=12.0)


# ═══════════════════════════════════════════════════
# SLIDE 27: 채널 전략 기회
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "05. 유통 채널 분석", "채널 전략 기회", "47")

# 가격 비교 차트 (채널 전략 시각화)
add_image_safe(slide, chart_path('chart_price_comparison.png'), 0.8, 1.3, width=12.0)
txbox(slide, 0.8, 5.5, 12, 0.3, "슈퍼세이브는 CJ온스타일 + 자사몰 중심에서 → 쿠팡/네이버/라이브커머스 확대가 핵심 과제", font_size=10, color=GRAY600)


# ═══════════════════════════════════════════════════
# SLIDE 28: 06 DIVIDER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_divider(slide, "06", "인사이트 & 기회", "Insights & Opportunities — SWOT, 핵심 기회, 전략 방향", "48")


# ═══════════════════════════════════════════════════
# SLIDE 29: 핵심 기회 포인트
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "06. 인사이트 & 기회", "핵심 기회 포인트", "49")

# 핵심 기회 포인트 차트
add_image_safe(slide, chart_path('chart_opportunities.png'), 0.8, 1.2, width=12.0)


# ═══════════════════════════════════════════════════
# SLIDE 30: 전략 방향 종합
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "06. 인사이트 & 기회", "전략 방향 종합", "50")

# 전략 로드맵 차트
add_image_safe(slide, chart_path('chart_roadmap.png'), 0.5, 1.2, width=12.5)


# ═══════════════════════════════════════════════════
# SLIDE 31: 07 DIVIDER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_divider(slide, "07", "해외 시장 분석", "Global Market — 글로벌 허벌 시장, 벤치마킹 브랜드, 진출 로드맵", "51")


# ═══════════════════════════════════════════════════
# SLIDES 32-34: 해외 시장 (요약)
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "07. 해외 시장 분석", "글로벌 허벌 시장 규모", "52")

global_kpis = [("$2,512억", "글로벌 허벌 시장 (2024)"), ("7.4%", "CAGR (2025-2030)"), ("$4,200억+", "2030 전망치"), ("아시아태평양", "최대 성장 지역")]
for i, (val, label) in enumerate(global_kpis):
    x = 0.8 + i * 3.0
    add_card_box(slide, x, 1.3, 2.7, 1.0)
    txbox(slide, x + 0.1, 1.35, 2.5, 0.5, val, font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txbox(slide, x + 0.1, 1.85, 2.5, 0.4, label, font_size=8, color=MUTED, align=PP_ALIGN.CENTER)

# 글로벌 시장 차트
add_image_safe(slide, chart_path('chart_global_market.png'), 0.8, 2.5, width=12.0)


slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "07. 해외 시장 분석", "해외 벤치마킹 브랜드", "53")

benchmarks = [
    ("Ricola", "스위스", "허벌 캔디/음료 글로벌 리더\n100년+ 전통, 50개국 유통"),
    ("Manuka Health", "뉴질랜드", "마누카꿀 프리미엄 브랜드\n원산지 스토리텔링 벤치마크"),
    ("Sambucol", "이스라엘→글로벌", "엘더베리 시럽 선도 브랜드\n면역력 포지셔닝 레퍼런스"),
    ("Traditional Medicinals", "미국", "허벌티 프리미엄 브랜드\nB Corp 인증, 지속가능성"),
    ("Comvita", "뉴질랜드", "마누카꿀 + 프로폴리스\n아시아 시장 진출 성공 사례"),
    ("Iberogast", "독일", "허벌 의약품 브랜드\n임상연구 기반 신뢰 구축"),
]
for i, (name, country, desc) in enumerate(benchmarks):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.4 + row * 2.6
    add_card_box(slide, x, y, 3.8, 2.2)
    txbox(slide, x + 0.15, y + 0.1, 3.5, 0.3, name, font_size=13, bold=True, color=BLACK)
    txbox(slide, x + 0.15, y + 0.45, 3.5, 0.2, country, font_size=9, bold=True, color=ACCENT)
    txbox(slide, x + 0.15, y + 0.75, 3.5, 1.2, desc, font_size=9, color=GRAY600)

    # 로고 영역 (추후 이미지 교체 가능)
    add_card_box(slide, x + 2.5, y + 0.1, 1.1, 0.8, GRAY100)


slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "07. 해외 시장 분석", "해외 진출 로드맵", "54")

phases = [
    ("Phase 1", "기반 구축 (1-6개월)", "• 브랜드 BI/패키징 글로벌화\n• 영문 자사몰 구축\n• 아마존 글로벌 셀링 입점\n• 해외 인증 준비 (FDA 등)", BLUE),
    ("Phase 2", "아시아 진출 (6-12개월)", "• 일본/대만/동남아 타겟\n• 현지 유통 파트너 확보\n• K-뷰티/K-푸드 연계 마케팅\n• 현지 인플루언서 협업", GREEN),
    ("Phase 3", "글로벌 확장 (12-24개월)", "• 미국/유럽 시장 진입\n• 오프라인 리테일 입점\n• 현지화 제품 개발\n• 글로벌 브랜드 캠페인", PURPLE),
]
for i, (phase, subtitle, items, color) in enumerate(phases):
    x = 0.8 + i * 4.2
    add_card_box(slide, x, 1.4, 3.8, 5.0)
    add_accent_line(slide, x, 1.4, 5.0, color)
    txbox(slide, x + 0.25, 1.55, 3.3, 0.3, phase, font_size=14, bold=True, color=color)
    txbox(slide, x + 0.25, 1.9, 3.3, 0.25, subtitle, font_size=10, bold=True, color=BLACK)
    txbox(slide, x + 0.25, 2.3, 3.3, 3.5, items, font_size=10, color=GRAY600)


# ═══════════════════════════════════════════════════
# SLIDE 35: CLOSING
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, DARK_BG, ACCENT)

txbox(slide, 3.5, 1.8, 6.5, 0.4, "THANK YOU", font_size=10, bold=True, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.CENTER)
txbox(slide, 3.5, 2.5, 6.5, 0.6, "Market Desk Research", font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(slide, 3.5, 3.3, 6.5, 0.4, "슈퍼세이브 — 곰보배추시럽", font_size=14, color=RGBColor(0xcc,0xcc,0xcc), align=PP_ALIGN.CENTER)
txbox(slide, 3.5, 3.8, 6.5, 0.4, "건강즙/허벌시럽 카테고리 · 국내외 시장 분석", font_size=11, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.CENTER)
txbox(slide, 3.5, 5.0, 6.5, 0.4, "Prepared by BRAND RISE · Confidential", font_size=9, color=RGBColor(0x77,0x77,0x77), align=PP_ALIGN.CENTER)
txbox(slide, 3.5, 5.4, 6.5, 0.3, "2026.04", font_size=9, color=RGBColor(0x77,0x77,0x77), align=PP_ALIGN.CENTER)


# ── 저장 ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "supersave-market-research.pptx")
prs.save(output_path)
print(f"✅ PPT 저장 완료: {output_path}")
print(f"   총 {len(prs.slides)} 슬라이드")
